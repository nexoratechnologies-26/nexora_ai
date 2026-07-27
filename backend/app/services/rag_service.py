import os
import re
import math
from typing import List, Dict, Any, Optional
import pypdf
import docx
import pptx
import chromadb
from chromadb.config import Settings as ChromaSettings
from app.core.config import settings

class RAGService:
    def __init__(self):
        self.chroma_client = chromadb.PersistentClient(
            path=settings.CHROMA_DB_DIR,
            settings=ChromaSettings(allow_reset=True)
        )

    def extract_text_from_file(self, file_path: str, file_type: str) -> str:
        """
        Extract raw text content based on file extension.
        """
        text = ""
        file_type = file_type.lower()
        
        if file_type == "txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
                
        elif file_type == "pdf":
            with open(file_path, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                        
        elif file_type == "docx":
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                if para.text:
                    text += para.text + "\n"
                    
        elif file_type == "pptx":
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
        else:
            raise ValueError(f"Unsupported file type for extraction: {file_type}")
            
        return text

    def chunk_text(self, text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
        """
        Split text into overlapping chunks.
        """
        # Normalize whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        if not text:
            return []
            
        words = text.split(' ')
        chunks = []
        
        # Simple word-based chunking for portability and performance
        step = chunk_size - chunk_overlap
        if step <= 0:
            step = chunk_size // 2
            
        i = 0
        while i < len(words):
            chunk_words = words[i:i + chunk_size]
            chunks.append(" ".join(chunk_words))
            i += step
            
        return chunks

    def _get_embedding(self, text: str, api_keys: Dict[str, str] = None) -> List[float]:
        """
        Compute an embedding vector for the text. 
        Uses OpenAI or Gemini embedding APIs if keys exist, and falls back to 
        a deterministic local pseudo-embedding (frequency/hash vector) 
        if offline or missing keys.
        """
        api_keys = api_keys or {}
        
        # 1. Try OpenAI Embedding if configured
        openai_key = api_keys.get("openai") or settings.OPENAI_API_KEY
        if openai_key:
            try:
                import openai
                client = openai.OpenAI(api_key=openai_key)
                response = client.embeddings.create(
                    input=[text],
                    model="text-embedding-3-small"
                )
                return response.data[0].embedding
            except Exception:
                pass # Fallback
                
        # 2. Try Gemini Embedding if configured
        gemini_key = api_keys.get("gemini") or settings.GEMINI_API_KEY
        if gemini_key:
            try:
                import google.generativeai as genai
                genai.configure(api_key=gemini_key)
                response = genai.embed_content(
                    model="models/embedding-001",
                    content=text,
                    task_type="retrieval_document"
                )
                return response["embedding"]
            except Exception:
                pass # Fallback

        # 3. Deterministic local fallback: Term Frequency Hash Vector (1536 dims)
        # Allows full offline testing and operations without API keys
        vector = [0.0] * 1536
        words = re.findall(r'\w+', text.lower())
        if not words:
            return vector
            
        for w in words:
            # Hash word to an index in 1536 dims
            idx = abs(hash(w)) % 1536
            vector[idx] += 1.0
            
        # L2 normalize
        magnitude = math.sqrt(sum(val * val for val in vector))
        if magnitude > 0:
            vector = [val / magnitude for val in vector]
            
        return vector

    async def ingest_document(
        self,
        document_id: str,
        user_id: str,
        file_path: str,
        file_type: str,
        api_keys: Dict[str, str] = None
    ) -> str:
        """
        Ingest a file: Extract text -> chunk -> generate embeddings -> store in ChromaDB.
        Returns the collection name.
        """
        text = self.extract_text_from_file(file_path, file_type)
        chunks = self.chunk_text(text)
        
        collection_name = f"user_{user_id.replace('-', '_')}"
        collection = self.chroma_client.get_or_create_collection(
            name=collection_name
        )
        
        if not chunks:
            return collection_name
            
        embeddings = [self._get_embedding(chunk, api_keys) for chunk in chunks]
        ids = [f"{document_id}_chunk_{i}" for i in range(len(chunks))]
        metadatas = [{"document_id": document_id, "user_id": user_id, "chunk_index": i} for i in range(len(chunks))]
        
        # ChromaDB accepts lists of document chunks, embeddings, metadatas, and ids
        collection.add(
            embeddings=embeddings,
            documents=chunks,
            ids=ids,
            metadatas=metadatas
        )
        
        return collection_name

    async def delete_document(self, document_id: str, user_id: str):
        """
        Delete all vectorized chunks associated with a document ID from ChromaDB.
        """
        collection_name = f"user_{user_id.replace('-', '_')}"
        try:
            collection = self.chroma_client.get_collection(name=collection_name)
            # Delete by metadata filter
            collection.delete(where={"document_id": document_id})
        except Exception:
            pass # Collection may not exist yet or was already deleted

    async def query_relevant_chunks(
        self,
        user_id: str,
        query: str,
        limit: int = 5,
        api_keys: Dict[str, str] = None
    ) -> List[Dict[str, Any]]:
        """
        Query user's vectorized collection for chunks relevant to the user request.
        """
        collection_name = f"user_{user_id.replace('-', '_')}"
        try:
            collection = self.chroma_client.get_collection(name=collection_name)
        except Exception:
            return [] # No collection created yet
            
        query_vector = self._get_embedding(query, api_keys)
        
        results = collection.query(
            query_embeddings=[query_vector],
            n_results=limit
        )
        
        chunks = []
        if results and "documents" in results and results["documents"]:
            docs = results["documents"][0]
            ids = results["ids"][0]
            metadatas = results["metadatas"][0]
            # Chroma DB distances are optional
            distances = results.get("distances", [[]])[0]
            
            for i in range(len(docs)):
                chunks.append({
                    "id": ids[i],
                    "text": docs[i],
                    "metadata": metadatas[i],
                    "score": float(distances[i]) if i < len(distances) else 1.0
                })
                
        return chunks
